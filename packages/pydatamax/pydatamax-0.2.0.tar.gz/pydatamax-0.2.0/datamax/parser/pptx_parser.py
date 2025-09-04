import builtins
import sys
import threading
from typing import Union

from loguru import logger

from datamax.parser.base import BaseLife, MarkdownOutputVo
from datamax.utils.lifecycle_types import LifeType

# Global variables for safe pptx import
_pptx_imported = False
_pptx_import_error = None
_pptx_import_lock = threading.Lock()
_pptx_presentation = None


def _safe_import_pptx():
    """Safely import pptx library to avoid conflicts with UNO (thread-safe)"""
    global _pptx_imported, _pptx_import_error, _pptx_presentation
    
    # Quick check to avoid unnecessary lock acquisition
    if _pptx_imported:
        return True
    
    with _pptx_import_lock:
        # Double-checked locking pattern
        if _pptx_imported:
            return True
        
        try:
            # Temporarily disable UNO's import hook
            original_import = None
            uno_module = sys.modules.get('uno')
            if uno_module and hasattr(uno_module, '_uno_import'):
                # Save the original import function
                original_import = builtins.__import__
                # Temporarily restore Python's original import
                builtins.__import__ = uno_module._builtin_import
            
            try:
                # Safely import pptx
                from pptx import Presentation
                _pptx_presentation = Presentation
                _pptx_imported = True
                logger.info("✅ PPTX module imported successfully")
                return True
            finally:
                # Restore UNO's import hook
                if original_import is not None:
                    builtins.__import__ = original_import
                    
        except ImportError as e:
            _pptx_import_error = e
            logger.error(f"❌ PPTX module import failed: {str(e)}")
            return False


def ensure_pptx_imported():
    """Ensure PPTX is imported, suitable for scenarios requiring early import"""
    if not _safe_import_pptx():
        raise ImportError(
            f"python-pptx library is not installed or cannot be imported. Error: {_pptx_import_error}\n"
            "Please install python-pptx: pip install python-pptx"
        )


class PptxParser(BaseLife):
    def __init__(self, file_path: str | list, domain: str = "Technology"):
        super().__init__(domain=domain)
        self.file_path = file_path

    @staticmethod
    def read_ppt_file(file_path: str):
        try:
            # Use safe import mechanism
            ensure_pptx_imported()
            
            content = ""
            prs = _pptx_presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        content += shape.text + "\n"
            return content
        except Exception:
            raise

    def parse(self, file_path: str) -> MarkdownOutputVo:
        # —— Lifecycle: Start processing PPTX —— #
        lc_start = self.generate_lifecycle(
            source_file=file_path,
            domain=self.domain,
            usage_purpose="Documentation",
            life_type=LifeType.DATA_PROCESSING,
        )
        logger.debug("⚙️ DATA_PROCESSING lifecycle generated")

        try:
            extension = self.get_file_extension(file_path)
            content = self.read_ppt_file(file_path=file_path)
            mk_content = content

            # —— Lifecycle: Processing completed —— #
            lc_end = self.generate_lifecycle(
                source_file=file_path,
                domain=self.domain,
                usage_purpose="Documentation",
                life_type=LifeType.DATA_PROCESSED,
            )
            logger.debug("⚙️ DATA_PROCESSED lifecycle generated")

            output_vo = MarkdownOutputVo(extension, mk_content)
            output_vo.add_lifecycle(lc_start)
            output_vo.add_lifecycle(lc_end)
            return output_vo.to_dict()

        except Exception as e:
            # —— Lifecycle: Processing failed —— #
            lc_fail = self.generate_lifecycle(
                source_file=file_path,
                domain=self.domain,
                usage_purpose="Documentation",
                life_type=LifeType.DATA_PROCESS_FAILED,
            )
            logger.debug("⚙️ DATA_PROCESS_FAILED lifecycle generated")

            raise Exception(
                {
                    "error": str(e),
                    "file_path": file_path,
                    "lifecycle": [lc_fail.to_dict()],
                }
            )
