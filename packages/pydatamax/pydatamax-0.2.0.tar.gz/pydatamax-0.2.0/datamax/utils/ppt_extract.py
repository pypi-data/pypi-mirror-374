import os
from functools import lru_cache
from pathlib import Path

from loguru import logger
from PIL.Image import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.base import BaseShape as Shape
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.group import GroupShape
from pptx.shapes.picture import Picture
from pptx.slide import Slide
from pptx.table import Table, _Cell, _Row


class PPtExtractor:
    @lru_cache(maxsize=128)
    def generate_img_path(self, id: str, img_name: str) -> str:
        if not isinstance(id, str):
            raise ValueError("id must be a string")
        if not isinstance(img_name, str):
            raise ValueError("img_name must be a string")
        return f"media/{id}/{img_name}"

    def handle_shape(
        self,
        shape: Shape,
        content_list: list[dict[str, str]],
        media_dir: Path,
        img_map: dict[Path, str],
        id: str,
        skip_image: bool,
    ):
        if not isinstance(shape, Shape):
            raise ValueError("Invalid shape object")
        if not isinstance(content_list, list):
            raise ValueError("content_list must be a list")
        if not isinstance(media_dir, Path):
            raise ValueError("media_dir must be a Path object")
        if not isinstance(img_map, dict):
            raise ValueError("img_map must be a dictionary")
        if not isinstance(id, str):
            raise ValueError("id must be a string")
        if not isinstance(skip_image, bool):
            raise ValueError("skip_image must be a boolean")

        try:
            shape_type = shape.shape_type
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    content_list.append(
                        {
                            "type": "text",
                            "data": paragraph.text + "\n",
                        }
                    )
            elif shape_type == MSO_SHAPE_TYPE.PICTURE and not skip_image:
                shape: Picture
                image: Image = shape.image
                image_bytes = image.blob
                img_path = media_dir.resolve().joinpath(
                    f"pic-{len(img_map)}.{image.ext}"
                )
                if not media_dir.exists():
                    media_dir.mkdir(parents=True, exist_ok=True)
                if not os.access(media_dir, os.W_OK):
                    raise PermissionError(f"Cannot write to directory: {media_dir}")
                img_s3_path = self.generate_img_path(id, img_path.name)
                img_map[img_path] = img_s3_path
                content_list.append({"type": "image", "data": img_s3_path})
                with open(img_path, "wb") as file:
                    file.write(image_bytes)
            elif shape_type == MSO_SHAPE_TYPE.TABLE:
                shape: GraphicFrame
                table: Table = shape.table
                md = "\n"
                for row_no, row in enumerate(table.rows):
                    row: _Row
                    md += "|"
                    if row_no == 1:
                        for col in row.cells:
                            md += "---|"
                        md += "\n|"
                    for col in row.cells:
                        cell: _Cell = col
                        md += (
                            " " + cell.text.replace("\r", " ").replace("\n", " ") + " |"
                        )
                    md += "\n"
                md += "\n"
                content_list.append({"type": "md", "data": md})
            elif shape_type == MSO_SHAPE_TYPE.GROUP:
                shape: GroupShape
                for sub_shape in shape.shapes:
                    self.handle_shape(
                        sub_shape, content_list, media_dir, img_map, id, skip_image
                    )
            else:
                logger.info(f"Unknown shape type: {shape_type}, {type(shape)}")
        except PermissionError as pe:
            logger.error(f"Permission error: {pe}")
        except OSError as ie:
            logger.error(f"IO error: {ie}")
        except Exception as e:
            logger.error(f"Error handling shape: {e}")

    def extract(
        self,
        presentation_source: Path,
        id: str,
        dir: Path,
        media_dir: Path,
        skip_image: bool,
    ):
        if not isinstance(presentation_source, Path):
            raise ValueError("presentation_source must be a Path object")
        if not isinstance(id, str):
            raise ValueError("id must be a string")
        if not isinstance(dir, Path):
            raise ValueError("dir must be a Path object")
        if not isinstance(media_dir, Path):
            raise ValueError("media_dir must be a Path object")
        if not isinstance(skip_image, bool):
            raise ValueError("skip_image must be a boolean")

        pages = []
        img_map = {}

        try:
            presentation: Presentation = Presentation(presentation_source)
            for page_no, slide in enumerate(presentation.slides):
                slide: Slide
                page = {"page_no": page_no, "content_list": []}
                for shape in slide.shapes:
                    self.handle_shape(
                        shape, page["content_list"], media_dir, img_map, id, skip_image
                    )
                pages.append(page)
        except FileNotFoundError as fnfe:
            logger.error(f"File not found: {fnfe}")
        except PermissionError as pe:
            logger.error(f"Permission error: {pe}")
        except OSError as ie:
            logger.error(f"IO error: {ie}")
        except Exception as e:
            logger.error(f"Error extracting presentation: {e}")

        return pages

    def run(self, id: str, file_path: Path, skip_image: bool = False):
        if not isinstance(id, str):
            raise ValueError("id must be a string")
        if not isinstance(file_path, Path):
            raise ValueError("file_path must be a Path object")
        if not isinstance(skip_image, bool):
            raise ValueError("skip_image must be a boolean")

        media_dir = Path("media").resolve()
        return self.extract(file_path, id, Path("."), media_dir, skip_image)
