"""
Advanced PDF to PowerPoint Converter

This module provides high-fidelity conversion from PDF documents to PowerPoint presentations,
preserving layouts, images, text formatting, and vector graphics.
"""

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.dml.color import RGBColor
from PIL import Image
import io
import os
import tempfile
import shutil
import math
import logging
from typing import List, Dict, Any, Tuple, Optional

# Configure logging
logger = logging.getLogger(__name__)

# Conversion factor from PDF points (1/72 inch) to EMUs (1/914400 inch)
POINTS_TO_EMU = 914400 / 72
WHITE_RGB = RGBColor(255, 255, 255)


class AdvancedPDFToPowerPointConverter:
    """
    Advanced PDF to PowerPoint converter that preserves high fidelity layouts.

    This converter extracts and reconstructs text, images, vector graphics, and
    formatting from PDF documents into PowerPoint presentations.
    """

    def __init__(self):
        """Initialize the converter with temporary directory for image processing."""
        self.temp_dir = tempfile.mkdtemp()
        self.slides_created = 0
        logger.debug(f"Initialized converter with temp directory: {self.temp_dir}")

    def _cleanup_temp_files(self):
        """Clean up temporary files created during conversion."""
        try:
            if os.path.exists(self.temp_dir):
                shutil.rmtree(self.temp_dir)
                logger.debug(f"Cleaned up temp directory: {self.temp_dir}")
        except Exception as e:
            logger.warning(f"Error cleaning up temp files: {e}")

    def _convert_color_tuple(
        self, color: Optional[Tuple[float, ...]]
    ) -> Optional[RGBColor]:
        """Convert color tuple from PDF to PowerPoint RGBColor."""
        if not color or len(color) < 3:
            return None
        return RGBColor(int(color[0] * 255), int(color[1] * 255), int(color[2] * 255))

    def _convert_srgb_color(self, srgb_int: Optional[int]) -> RGBColor:
        """Convert sRGB integer color to PowerPoint RGBColor."""
        if srgb_int is None:
            return RGBColor(0, 0, 0)
        return RGBColor((srgb_int >> 16) & 255, (srgb_int >> 8) & 255, srgb_int & 255)

    def _extract_page_elements(
        self, pdf_doc: fitz.Document, start_page: int, end_page: int
    ) -> List[Dict[str, Any]]:
        """
        Extract all elements (text, images, vector graphics) from PDF pages.

        Args:
            pdf_doc: The PDF document object
            start_page: Starting page index (0-based)
            end_page: Ending page index (0-based)

        Returns:
            List of dictionaries containing page data and elements
        """
        all_pages_data = []

        for page_num in range(start_page, end_page + 1):
            page = pdf_doc.load_page(page_num)
            page_data = {"page_number": page_num + 1, "elements": []}
            logger.debug(f"Extracting elements from page {page_num + 1}")

            # Extract Vector Graphics
            for path in page.get_drawings():
                for item in path["items"]:
                    element = None
                    if item[0] == "l":  # Line
                        element = {"type": "line", "p1": item[1], "p2": item[2]}
                    elif item[0] == "re":  # Rectangle
                        element = {"type": "rect", "bbox": fitz.Rect(item[1])}

                    if element:
                        element.update(
                            {
                                "line_color": self._convert_color_tuple(path["color"]),
                                "fill_color": self._convert_color_tuple(path["fill"]),
                                "line_width": path["width"],
                            }
                        )
                        page_data["elements"].append(element)

            # Extract Images
            for img_index, img in enumerate(page.get_images(full=True)):
                xref, smask_xref = img[0], img[1]
                try:
                    img_info = pdf_doc.extract_image(xref)
                    img_bytes, img_ext = img_info["image"], img_info["ext"]

                    # Handle images with transparency masks
                    if smask_xref > 0:
                        mask_info = pdf_doc.extract_image(smask_xref)
                        with Image.open(
                            io.BytesIO(img_info["image"])
                        ) as base_img, Image.open(
                            io.BytesIO(mask_info["image"])
                        ).convert(
                            "L"
                        ) as mask_img:
                            base_img = base_img.convert("RGBA")
                            base_img.putalpha(mask_img)
                            img_ext = "png"
                            with io.BytesIO() as final_img_bytes:
                                base_img.save(final_img_bytes, format="PNG")
                                img_bytes = final_img_bytes.getvalue()

                    img_bbox = page.get_image_bbox(img)
                    if img_bbox and not img_bbox.is_empty:
                        path = os.path.join(
                            self.temp_dir, f"p{page_num}_i{img_index}.{img_ext}"
                        )
                        with open(path, "wb") as f:
                            f.write(img_bytes)
                        page_data["elements"].append(
                            {"type": "image", "bbox": img_bbox, "path": path}
                        )

                except Exception as e:
                    logger.warning(
                        f"Could not process image {img_index} on page {page_num + 1}: {e}"
                    )

            # Extract Text (Per-Line)
            for block in page.get_text("dict", flags=fitz.TEXTFLAGS_TEXT)["blocks"]:
                if block["type"] == 0:  # Text block
                    for line in block["lines"]:
                        line_bbox = fitz.Rect(line["bbox"])
                        if line_bbox.is_empty or not line["spans"]:
                            continue

                        element = {"type": "textbox", "bbox": line_bbox, "spans": []}
                        for span in line["spans"]:
                            flags = span["flags"]
                            element["spans"].append(
                                {
                                    "text": span["text"],
                                    "font": span["font"],
                                    "size": span["size"],
                                    "color": self._convert_srgb_color(span["color"]),
                                    "bold": bool(flags & (1 << 4)),
                                    "italic": bool(flags & (1 << 1)),
                                }
                            )
                        page_data["elements"].append(element)

            all_pages_data.append(page_data)

        return all_pages_data

    def _filter_knockout_rects(self, elements: List[Dict]) -> List[int]:
        """
        Filter out white rectangles that are likely knockout elements behind text.

        Args:
            elements: List of page elements

        Returns:
            List of indices of elements to skip
        """
        text_bboxes = [el["bbox"] for el in elements if el["type"] == "textbox"]
        rect_indices_to_skip = []

        for i, el in enumerate(elements):
            if (
                el["type"] == "rect"
                and el["fill_color"] == WHITE_RGB
                and not el["line_color"]
            ):

                for text_bbox in text_bboxes:
                    if el["bbox"].contains(text_bbox) and (
                        el["bbox"].get_area() / text_bbox.get_area() < 2.0
                    ):
                        rect_indices_to_skip.append(i)
                        break

        return rect_indices_to_skip

    def _create_slide_from_page(self, prs: Presentation, page_data: Dict):
        """
        Create a PowerPoint slide from extracted page elements.

        Args:
            prs: PowerPoint presentation object
            page_data: Dictionary containing page elements
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        elements = page_data["elements"]
        knockout_indices = self._filter_knockout_rects(elements)

        for i, el in enumerate(elements):
            if i in knockout_indices:
                continue

            try:
                if el["type"] == "rect":
                    self._add_rectangle(slide, el)
                elif el["type"] == "line":
                    self._add_line(slide, el)
                elif el["type"] == "image":
                    self._add_image(slide, el)
                elif el["type"] == "textbox":
                    self._add_textbox(slide, el)

            except Exception as e:
                logger.warning(f"Could not add element of type {el.get('type')}: {e}")

    def _add_rectangle(self, slide, element):
        """Add a rectangle shape to the slide."""
        bbox = element["bbox"]
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(bbox.x0 * POINTS_TO_EMU),
            Emu(bbox.y0 * POINTS_TO_EMU),
            Emu(bbox.width * POINTS_TO_EMU),
            Emu(bbox.height * POINTS_TO_EMU),
        )

        # Set fill color
        if element["fill_color"]:
            shape.fill.solid()
            shape.fill.fore_color.rgb = element["fill_color"]
        else:
            shape.fill.background()

        # Set line color and width
        if element["line_color"]:
            shape.line.color.rgb = element["line_color"]
            shape.line.width = Emu(max(element["line_width"], 0.5) * POINTS_TO_EMU)
        else:
            shape.line.fill.background()

    def _add_line(self, slide, element):
        """Add a line shape to the slide."""
        p1, p2 = element["p1"], element["p2"]
        x1, y1 = p1.x * POINTS_TO_EMU, p1.y * POINTS_TO_EMU
        x2, y2 = p2.x * POINTS_TO_EMU, p2.y * POINTS_TO_EMU

        width = math.hypot(x2 - x1, y2 - y1)
        if width < 1:
            return

        height = Emu(max(element["line_width"], 0.5) * POINTS_TO_EMU)
        angle_deg = math.degrees(math.atan2(y2 - y1, x2 - x1))
        left = ((x1 + x2) / 2) - (width / 2)
        top = ((y1 + y2) / 2) - (height / 2)

        line_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(height)
        )
        line_shape.rotation = angle_deg

        if element["line_color"]:
            line_shape.fill.solid()
            line_shape.fill.fore_color.rgb = element["line_color"]
            line_shape.line.fill.background()
        else:
            line_shape.fill.background()

    def _add_image(self, slide, element):
        """Add an image to the slide."""
        bbox = element["bbox"]
        slide.shapes.add_picture(
            element["path"],
            Emu(bbox.x0 * POINTS_TO_EMU),
            Emu(bbox.y0 * POINTS_TO_EMU),
            Emu(bbox.width * POINTS_TO_EMU),
            Emu(bbox.height * POINTS_TO_EMU),
        )

    def _add_textbox(self, slide, element):
        """Add a textbox with formatted text to the slide."""
        bbox = element["bbox"]
        textbox = slide.shapes.add_textbox(
            Emu(bbox.x0 * POINTS_TO_EMU),
            Emu(bbox.y0 * POINTS_TO_EMU),
            Emu(bbox.width * POINTS_TO_EMU),
            Emu(bbox.height * POINTS_TO_EMU),
        )

        tf = textbox.text_frame
        tf.clear()
        tf.margin_bottom = tf.margin_top = tf.margin_left = tf.margin_right = Emu(0)
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP

        p = tf.paragraphs[0]
        p.space_before = p.space_after = Pt(0)

        for span in element["spans"]:
            run = p.add_run()
            run.text = span["text"]
            font = run.font
            font.name = span["font"]
            font.size = Pt(span["size"])
            font.color.rgb = span["color"]
            font.bold = span["bold"]
            font.italic = span["italic"]

    def convert(
        self,
        pdf_path: str,
        output_path: str,
        page_range: Optional[Tuple[int, int]] = None,
    ) -> bool:
        """
        Convert a PDF document to PowerPoint presentation.

        Args:
            pdf_path: Path to the input PDF file
            output_path: Path for the output PowerPoint file
            page_range: Optional tuple (start_page, end_page) for partial conversion (1-based)

        Returns:
            bool: True if conversion successful, False otherwise

        Raises:
            FileNotFoundError: If the input PDF file doesn't exist
            ValueError: If page range is invalid
        """
        if not os.path.exists(pdf_path):
            raise FileNotFoundError(f"Input PDF not found at '{pdf_path}'")

        logger.info("Starting PDF to PowerPoint conversion...")

        try:
            pdf_doc = fitz.open(pdf_path)
            doc_page_count = len(pdf_doc)
            logger.info(f"Opened PDF with {doc_page_count} pages")

            # Validate and set page range
            start_page, end_page = 0, doc_page_count - 1

            if page_range:
                user_start, user_end = page_range

                if not (isinstance(user_start, int) and isinstance(user_end, int)):
                    raise ValueError(
                        f"Page range values must be integers. Got {page_range}"
                    )

                if user_start > user_end:
                    raise ValueError(
                        f"Start page ({user_start}) cannot be greater than end page ({user_end})"
                    )

                if not (
                    1 <= user_start <= doc_page_count
                    and 1 <= user_end <= doc_page_count
                ):
                    raise ValueError(
                        f"Page range {page_range} is out of bounds. Document has {doc_page_count} pages"
                    )

                start_page, end_page = (
                    user_start - 1,
                    user_end - 1,
                )  # Convert to 0-based index
                logger.info(
                    f"Processing user-defined page range: {user_start} to {user_end}"
                )
            else:
                logger.info(f"Processing all {doc_page_count} pages")

            # Get PDF dimensions from first page
            first_page = pdf_doc.load_page(start_page)
            pdf_width_pt, pdf_height_pt = first_page.rect.width, first_page.rect.height
            logger.debug(f"PDF dimensions: {pdf_width_pt}x{pdf_height_pt} points")

            # Create PowerPoint presentation with custom dimensions
            prs = Presentation()
            prs.slide_width = Emu(pdf_width_pt * POINTS_TO_EMU)
            prs.slide_height = Emu(pdf_height_pt * POINTS_TO_EMU)
            logger.debug("Created PowerPoint presentation with custom dimensions")

            # Extract and convert pages
            logger.info("Extracting and reconstructing page elements...")
            all_pages_data = self._extract_page_elements(pdf_doc, start_page, end_page)

            for i, page_data in enumerate(all_pages_data):
                logger.debug(f"Processing page {start_page + i + 1}/{end_page + 1}")
                self._create_slide_from_page(prs, page_data)

            pdf_doc.close()
            prs.save(output_path)
            self.slides_created = len(prs.slides)

            logger.info(f"Conversion successful! Output saved to: {output_path}")
            logger.info(f"Created {self.slides_created} slides")
            return True

        except Exception as e:
            logger.error(f"Conversion failed: {e}")
            raise
        finally:
            self._cleanup_temp_files()

    def __enter__(self):
        """Context manager entry."""
        return self

    def __exit__(self, exc_type, exc_val, exc_tb):
        """Context manager exit with cleanup."""
        self._cleanup_temp_files()
