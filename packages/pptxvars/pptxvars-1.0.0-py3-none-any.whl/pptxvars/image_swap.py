from pathlib import Path
import sys
import yaml
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

__all__ = ["swap_frames_from_imgs"]


def _collect_frames(slide):
    frames = []
    def add_item(shp, left_abs, top_abs, parent_group=None):
        frames.append({
            "shape": shp,
            "left": left_abs,
            "top": top_abs,
            "width": shp.width,
            "height": shp.height,
            "rotation": getattr(shp, "rotation", 0),
            "parent_group": parent_group,
        })
    for shp in slide.shapes:
        if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
            g = shp
            for child in g.shapes:
                if "frame" in (child.name or "").lower():
                    add_item(child, g.left + child.left, g.top + child.top, parent_group=g)
        else:
            if "frame" in (shp.name or "").lower():
                add_item(shp, shp.left, shp.top, parent_group=None)
    frames.sort(key=lambda x: (int(x["top"]), int(x["left"])))
    return frames


def _add_picture_fit(slide, img_path, left, top, width, height, rotation):
    pic = slide.shapes.add_picture(str(img_path), left, top, width=width)
    if pic.height > height:
        pic.height = height
    pic.left = left + (width - pic.width) // 2
    pic.top  = top  + (height - pic.height) // 2
    pic.rotation = rotation
    return pic


def _swap_slide_shape(slide, shp, img_path, left, top, width, height, rotation):
    tree = slide.shapes._spTree
    z = list(tree).index(shp._element)
    tree.remove(shp._element)
    pic = _add_picture_fit(slide, img_path, left, top, width, height, rotation)
    tree.remove(pic._element); tree.insert(z, pic._element)


def _swap_group_child(slide, group, child, img_path, left_abs, top_abs, width, height, rotation):
    group.shapes._spTree.remove(child._element)
    pic = _add_picture_fit(slide, img_path, left_abs, top_abs, width, height, rotation)
    st = slide.shapes._spTree
    gi = list(st).index(group._element)
    st.remove(pic._element); st.insert(gi + 1, pic._element)


def _load_imgs(yml_path: Path):
    data = yaml.safe_load(yml_path.read_text(encoding="utf-8")) or {}
    slides = data.get("slides") or []
    return [(int(item["index"]), [Path(p) for p in item.get("images", [])]) for item in slides]


def swap_frames_from_imgs(pptx_in: Path, imgs_yml: Path, pptx_out: Path):
    prs = Presentation(pptx_in)
    plan = _load_imgs(imgs_yml)
    if not plan:
        sys.exit("YAML must define: slides: - index: <int>  images: [paths...]")
    base = imgs_yml.parent

    for idx, imgs in plan:
        if not (1 <= idx <= len(prs.slides)):
            print(f"WARN: slide {idx} missing", file=sys.stderr); continue
        slide = prs.slides[idx - 1]
        frames = _collect_frames(slide)
        if not frames:
            print(f"WARN: slide {idx}: no frames (name contains 'Frame')", file=sys.stderr); continue

        for frame, rel_img in zip(frames, imgs):
            img_path = (base / rel_img).resolve()
            if not img_path.exists():
                print(f"WARN: slide {idx}: image not found: {img_path}", file=sys.stderr); continue
            if frame["parent_group"] is None:
                _swap_slide_shape(slide, frame["shape"], img_path,
                                  frame["left"], frame["top"], frame["width"], frame["height"], frame["rotation"])
            else:
                _swap_group_child(slide, frame["parent_group"], frame["shape"], img_path,
                                  frame["left"], frame["top"], frame["width"], frame["height"], frame["rotation"])

        if len(imgs) > len(frames):
            print(f"WARN: slide {idx}: {len(imgs)-len(frames)} images unused (only {len(frames)} frames)", file=sys.stderr)
        if len(frames) > len(imgs):
            print(f"NOTE: slide {idx}: {len(frames)-len(imgs)} frames left unchanged", file=sys.stderr)

    pptx_out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(pptx_out)
