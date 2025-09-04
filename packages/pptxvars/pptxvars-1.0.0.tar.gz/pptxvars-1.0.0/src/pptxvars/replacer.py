import re
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import yaml

__all__ = ["apply_vars", "load_vars", "format_outpath"]


def load_vars(yaml_path: Path) -> dict:
    data = yaml.safe_load(open(yaml_path, "r", encoding="utf-8")) or {}
    return {str(k).upper(): "" if v is None else str(v) for k, v in data.items()}


def build_regex(keys):
    inner = "|".join(re.escape(k) for k in keys)
    return re.compile(r"\{\{\s*(" + inner + r")\s*\}\}")


def paragraph_text_and_spans(p):
    full, spans, pos = [], [], 0
    for r in p.runs:
        t = r.text or ""
        full.append(t); spans.append((pos, pos+len(t))); pos += len(t)
    return "".join(full), spans


def replace_in_paragraph_keep_runs(p, rx, values):
    full, spans = paragraph_text_and_spans(p)
    if not full: return False
    matches = list(rx.finditer(full))
    if not matches: return False

    for m in reversed(matches):
        key = m.group(1).upper()
        repl = values.get(key, "")
        s, e = m.span()
        start_i = next(i for i,(a,b) in enumerate(spans) if a <= s < b)
        end_i   = next(i for i,(a,b) in enumerate(spans) if a <  b and e-1 < b)
        for i in range(start_i, end_i+1):
            run = p.runs[i]; a,b = spans[i]
            ls = max(0, s-a); le = min(b-a, e-a)
            before = run.text[:ls]; after = run.text[le:]
            if i == start_i:
                run.text = before + repl + after
                delta = len(run.text) - (b-a)
            else:
                run.text = before + after
                delta = len(run.text) - (b-a)
            spans[i] = (a, b + delta)
            for j in range(i+1, len(spans)):
                sa,sb = spans[j]; spans[j] = (sa+delta, sb+delta)
    return True


def walk_shapes(shapes, rx, values):
    changed = False
    for shp in list(shapes):
        if getattr(shp, "has_text_frame", False):
            for p in shp.text_frame.paragraphs:
                changed |= replace_in_paragraph_keep_runs(p, rx, values)
        if getattr(shp, "has_table", False):
            for row in shp.table.rows:
                for cell in row.cells:
                    if cell.text_frame:
                        for p in cell.text_frame.paragraphs:
                            changed |= replace_in_paragraph_keep_runs(p, rx, values)
        if shp.shape_type == MSO_SHAPE_TYPE.GROUP and hasattr(shp, "shapes"):
            changed |= walk_shapes(shp.shapes, rx, values)
    return changed


def apply_vars(pptx_in: Path, yaml_vars: Path, pptx_out: Path):
    vals = load_vars(yaml_vars)
    if not vals:
        raise SystemExit("No variables in YAML.")
    rx = build_regex(vals.keys())
    prs = Presentation(pptx_in)

    for m in prs.slide_masters:
        walk_shapes(m.shapes, rx, vals)
    for lay in prs.slide_layouts:
        walk_shapes(lay.shapes, rx, vals)
    for s in prs.slides:
        walk_shapes(s.shapes, rx, vals)

    pptx_out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(pptx_out)


def _safe_filename(s: str) -> str:
    s = s.strip().replace(" ", "_")
    return re.sub(r'[\\/*?:"<>|]+', "-", s)


def format_outpath(pattern: str, variables: dict, default_dir: Path) -> Path:
    # pattern like "output/{STEM}_{DATE}.pptx"
    class SafeDict(dict):
        def __missing__(self, k): return "{"+k+"}"
    filled = pattern.format_map(SafeDict(**variables))
    parts = Path(filled)
    filename = _safe_filename(parts.name)
    return (default_dir / parts.parent / filename).resolve()
