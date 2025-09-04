from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import os


def b_imgs2ppt(
        image_folder,
        output_ppt="output.pptx",
        base_image_path=None,
        sorted_key=lambda x: int(x.split(".")[0]),
        dpi=200
):
    '''
    将图片文件夹中的图片转换为 PPT 文件
    :param image_folder:
    :param output_ppt:
    :param base_image_path: 默认使用第一个图片作为基准图片
    :param sorted_key:
    :param dpi:
    :return:
    '''
    # 获取所有图片并排序
    images = sorted([
        file for file in os.listdir(image_folder)
        if file.lower().endswith((".png", ".jpg", ".jpeg", ".bmp", ".gif"))
    ])
    images = sorted(images, key=sorted_key)  # 按照数字大小排序

    prs = Presentation()
    # 获取基准图片的尺寸（像素）
    if base_image_path is None:
        base_image_path = os.path.join(image_folder, images[0])
    with Image.open(base_image_path) as img:
        width_px, height_px = img.size
    # 转换为英寸
    width_in = width_px / dpi
    height_in = height_px / dpi
    # 设置幻灯片尺寸（单位是英寸，需转成英寸再乘 Inches）
    prs.slide_width = Inches(width_in)
    prs.slide_height = Inches(height_in)

    # 遍历所有图片并插入到幻灯片中
    for img_name in images:
        img_path = os.path.join(image_folder, img_name)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    prs.save(output_ppt)
    print(f"PPT已保存到: {output_ppt}")
