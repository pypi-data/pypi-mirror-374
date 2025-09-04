from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import os


def b_imgs2ppt(
        image_folder,
        output_ppt="output.pptx",
        sorted_key=lambda x: int(x.split(".")[0]),
        dpi=96  # 改成96，更接近屏幕分辨率
):
    images = sorted([
        file for file in os.listdir(image_folder)
        if file.lower().endswith((".png", ".jpg", ".jpeg", ".bmp", ".gif"))
    ], key=sorted_key)

    prs = Presentation()

    for img_name in images:
        img_path = os.path.join(image_folder, img_name)

        with Image.open(img_path) as img:
            width_px, height_px = img.size

        # 转换为英寸
        width_in = max(width_px / dpi, 1)   # 最小1英寸
        height_in = max(height_px / dpi, 1) # 最小1英寸

        # 修改PPT页面尺寸
        prs.slide_width = Inches(width_in)
        prs.slide_height = Inches(height_in)

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    prs.save(output_ppt)
    print(f"PPT已保存到: {output_ppt}")
