from typing import Any, Dict, List, Optional, Tuple
from abc import ABC, abstractmethod
import tempfile
from pathlib import Path
from io import BytesIO
import markdown
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
from navconfig.logging import logging


class PowerPointFile(ABC):
    """
    PowerPointFile Interface

    Provides common functionality for working with PowerPoint presentations:
    - Template loading and layout analysis
    - Hybrid template/programmatic layout creation
    - Image and text handling
    - Markdown conversion
    - Style application
    - File operations
    """

    def __init__(self, *args, **kwargs):
        self.logger = logging.getLogger(
            'Flowtask.PowerPointFile'
        )
        self._layout_cache: Dict[int, Dict[str, Any]] = {}
        self._presentation: Optional[Presentation] = None
        self.text_styles: Dict[str, Any] = kwargs.get('text_styles', {
            'font_name': 'Calibri',
            'font_size': 12,
            'bold': False
        })
        self.image_size: Dict[str, float] = kwargs.get(
            'image_size', {
                'width': 4.0,
                'height': 3.0
            }
        )
        super().__init__(*args, **kwargs)

    def load_template(self, template_path: Path) -> bool:
        """Load PowerPoint template and analyze layouts."""
        try:
            self._presentation = Presentation(str(template_path))
            self.logger.info(f"Loaded PowerPoint template: {template_path}")
            return True
        except Exception as e:
            self.logger.error(f"Failed to load PowerPoint template: {e}")
            return False

    def clear_existing_slides(self):
        """Remove existing slides."""
        if not self._presentation:
            return
        while len(self._presentation.slides) > 0:
            rId = self._presentation.slides._sldIdLst[0].rId
            self._presentation.part.drop_rel(rId)
            del self._presentation.slides._sldIdLst[0]

    # Utility methods for text processing
    def convert_markdown_to_html(self, markdown_text: str) -> str:
        """Convert markdown to HTML."""
        try:
            html = markdown.markdown(
                markdown_text,
                extensions=['extra', 'tables']
            )
            return html
        except Exception as e:
            self.logger.error(f"Failed to convert markdown to HTML: {e}")
            return markdown_text

    def extract_text_from_html(self, html: str) -> str:
        """Extract clean text from HTML for PowerPoint."""
        try:
            soup = BeautifulSoup(html, 'html.parser')
            text_parts = []

            for element in soup.find_all(['p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'li', 'div']):
                text = element.get_text().strip()
                if text:
                    if element.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                        text = f"• {text}"
                    elif element.name == 'li':
                        text = f"  - {text}"
                    text_parts.append(text)

            return '\n'.join(text_parts)
        except Exception as e:
            self.logger.error(f"Failed to extract text from HTML: {e}")
            return html

    def save_presentation(self, output_filename, output_path: Path, override: bool = True):
        """Save the presentation to file."""
        if not self._presentation:
            raise RuntimeError("No presentation to save")
        output_filename = output_filename if output_filename.endswith('.pptx') else f"{output_filename}.pptx"
        path = output_path.joinpath(output_filename)

        if not override and path.exists():
            self.logger.warning(f"PowerPoint presentation already exists: {path}")
            return
        if override and path.exists():
            self.logger.info(f"Overriding existing PowerPoint presentation: {path}")
            path.unlink(missing_ok=True)

        self._presentation.save(str(path))
        self.logger.notice(
            f"PowerPoint presentation saved: {path}"
        )

    def _layout_catalog(self):
        """
        Scan all layouts and collect {layout_idx: {placeholder_name: placeholder_idx}}
        """
        catalog = {}
        for i, ly in enumerate(self._presentation.slide_layouts):
            names = {}
            for ph in getattr(ly, "placeholders", []):
                try:
                    nm = getattr(ph, "name", None)
                    idx = ph.placeholder_format.idx
                except Exception:
                    continue
                if nm:
                    names[nm] = idx
            catalog[i] = names
        return catalog

    def _choose_best_layout_index(self) -> int:
        """
        Choose the layout that contains the most placeholder names from column_mapping.
        Falls back to self.slide_layout, then 1, then 6.
        """
        wanted = {dest for dest in self.column_mapping.values() if isinstance(dest, str)}
        catalog = self._layout_catalog()
        best_idx, best_score = None, -1
        for idx, names in catalog.items():
            score = len(wanted.intersection(set(names.keys())))
            if score > best_score:
                best_idx, best_score = idx, score
        for fb in (getattr(self, 'slide_layout', 1), 1, 6):
            if best_idx is None and fb < len(self._presentation.slide_layouts):
                best_idx = fb
        return best_idx if best_idx is not None else 1

    def _safe_placeholder_idx(self, shape):
        try:
            return shape.placeholder_format.idx
        except Exception:
            return None

    def _find_placeholder_by_name(self, slide, placeholder_name: str):
        # Direct name match first
        for shp in slide.shapes:
            if getattr(shp, "name", None) == placeholder_name:
                return shp

        # Match by placeholder index from the layout
        ly = slide.slide_layout
        for lph in getattr(ly, "placeholders", []):
            if getattr(lph, "name", None) == placeholder_name:
                target_idx = self._safe_placeholder_idx(lph)
                if target_idx is None:
                    continue
                for shp in slide.shapes:
                    if self._safe_placeholder_idx(shp) == target_idx:
                        return shp
        return None

    def _ensure_placeholder(self, slide, section_name: str, autoplace_state: dict):
        """
        Return a shape for section_name.
        If missing, create it. Uses self.missing_sections[section_name] if present;
        otherwise autoplaces in a flowing vertical stack from top-left margin.
        """
        shp = self._find_placeholder_by_name(slide, section_name)
        if shp:
            # If off-canvas (e.g., Date placeholder moved out), center in footer
            if shp.left < 0 or shp.top < 0:
                sw, sh = self._presentation.slide_width, self._presentation.slide_height
                shp.left = Inches(0.5)
                shp.top = sh - Inches(0.6)
                shp.width = sw - Inches(1.0)
            return shp

        spec = getattr(self, "missing_sections", {}).get(section_name)

        if spec:
            # anchor-based position
            anchor = None
            if a := spec.get("anchor"):
                anchor = self._find_placeholder_by_name(slide, a)
            pos = (spec.get("position") or "below").lower()
            off = spec.get("offset_in", {})
            dx, dy = Inches(off.get("dx", 0.0)), Inches(off.get("dy", 0.0))
            sz = spec.get("size_in", {})
            width, height = Inches(sz.get("width", 3.0)), Inches(sz.get("height", 0.5))
            if anchor:
                if pos == "below":
                    left, top = anchor.left, anchor.top + anchor.height + dy
                elif pos == "right":
                    left, top = anchor.left + anchor.width + dx, anchor.top
                elif pos == "left":
                    left, top = anchor.left - width + dx, anchor.top
                else:  # above
                    left, top = anchor.left, anchor.top - height + dy
            else:
                left = Inches(spec.get("left_in", 0.5))
                top = Inches(spec.get("top_in", 0.5))
            box = slide.shapes.add_textbox(left, top, width, height)
            box.name = section_name
            self.logger.debug(f"Created missing section '{section_name}' at position ({left}, {top})")
            return box

        # Auto-place fallback: simple flowing stack
        left = autoplace_state.setdefault("left", Inches(0.5))
        top = autoplace_state.setdefault("top", Inches(0.5))
        width = Inches(4.0)
        height = Inches(0.45)
        box = slide.shapes.add_textbox(left, top, width, height)
        box.name = section_name
        # advance cursor for the next missing shape
        autoplace_state["top"] = top + height + Inches(0.1)
        self.logger.debug(f"Auto-placed section '{section_name}' at position ({left}, {top})")
        return box

    def _normalize_image(self, value):
        """
        Normalize various image formats to PIL.Image.
        Handles BytesIO, PIL.Image, file paths, numpy arrays, etc.
        """
        from PIL import Image as PILImage
        import pandas as pd

        try:
            if isinstance(value, PILImage.Image):
                return value
            elif isinstance(value, BytesIO):
                value.seek(0)  # Reset position
                img = PILImage.open(value)
                # Convert to RGB if needed for compatibility
                if img.mode in ('RGBA', 'LA', 'P'):
                    img = img.convert('RGB')
                return img
            elif isinstance(value, (bytes, bytearray)):
                img = PILImage.open(BytesIO(value))
                if img.mode in ('RGBA', 'LA', 'P'):
                    img = img.convert('RGB')
                return img
            elif isinstance(value, str) and Path(value).exists():
                img = PILImage.open(value)
                if img.mode in ('RGBA', 'LA', 'P'):
                    img = img.convert('RGB')
                return img
            elif hasattr(value, 'read'):  # File-like object
                img = PILImage.open(value)
                if img.mode in ('RGBA', 'LA', 'P'):
                    img = img.convert('RGB')
                return img
            elif "numpy" in str(type(value)):
                return PILImage.fromarray(value)
            elif pd.isna(value) or value is None:
                return None
            else:
                self.logger.warning(f"Unrecognized image format: {type(value)}")
                return None
        except Exception as e:
            self.logger.error(f"Failed to normalize image: {e}")
            return None

    def _format_value(self, col, val):
        """Enhanced value formatting for different data types and contexts."""
        import pandas as pd
        from datetime import datetime, date

        if pd.isna(val) or val is None:
            return getattr(self, 'reference_text', {}).get(col, f"[{col}]")

        # Handle boolean values with context-aware formatting
        if isinstance(val, bool):
            if any(word in col.lower() for word in ['compliant', 'status', 'pass', 'fail']):
                return "✓ Compliant" if val else "✗ Non-Compliant"
            elif any(word in col.lower() for word in ['active', 'enabled', 'valid']):
                return "✓ Yes" if val else "✗ No"
            return "True" if val else "False"

        # Handle datetime objects
        if isinstance(val, (pd.Timestamp, datetime, date)):
            if hasattr(val, 'strftime'):
                return val.strftime("%Y-%m-%d %H:%M") if hasattr(val, 'hour') else val.strftime("%Y-%m-%d")
            return str(val)

        # Handle numeric scores and percentages
        if isinstance(val, float):
            if any(word in col.lower() for word in ['score', 'rate', 'percent', 'pct']):
                # If value is between 0 and 1, treat as percentage
                if 0 <= val <= 1:
                    return f"{val:.1%}"
                else:
                    return f"{val:.1f}"
            # General float formatting
            return f"{val:.2f}" if val != int(val) else str(int(val))

        # Handle store IDs and similar identifiers
        if 'store_id' in col.lower() or 'id' in col.lower():
            return f"Store ID: {val}" if 'store' in col.lower() else f"ID: {val}"

        # Default string conversion
        return str(val)

    def _apply_text_styles(self, text_frame, section_name: str, custom_styles: dict = None):
        """Apply text styles with enhanced section-specific formatting."""
        try:
            # Enhanced section-specific styles
            default_styles = {
                'store_id_text': {'font_size': 16, 'bold': True, 'color': '003366'},
                'created_on_text': {'font_size': 11, 'bold': False, 'color': '666666'},
                'who_text': {'font_size': 12, 'bold': False, 'color': '666666'},
                'visitor_name': {'font_size': 12, 'bold': False, 'color': '666666'},
                'analysis_text': {'font_size': 10, 'bold': False, 'color': '333333'},
                'score_text': {'font_size': 14, 'bold': True, 'color': '008000'},
                'status_text': {'font_size': 12, 'bold': True, 'color': '008000'},
                'photo_taken_text': {'font_size': 9, 'bold': False, 'color': '808080'},
                'compliance_analysis_markdown': {'font_size': 10, 'bold': False, 'color': '333333'}
            }

            # Merge custom styles if provided
            if custom_styles:
                default_styles.update(custom_styles)

            styles = default_styles.get(section_name, {'font_size': 12, 'bold': False})

            for paragraph in text_frame.paragraphs:
                font = paragraph.font
                font.name = "Calibri"

                # Apply font size
                if 'font_size' in styles:
                    font.size = Pt(int(styles['font_size']))

                font.bold = styles.get('bold', False)

                # Apply color
                if 'color' in styles:
                    color_hex = styles['color']
                    if len(color_hex) == 6:  # Ensure valid hex color
                        try:
                            r = int(color_hex[0:2], 16)
                            g = int(color_hex[2:4], 16)
                            b = int(color_hex[4:6], 16)
                            font.color.rgb = RGBColor(r, g, b)
                        except ValueError:
                            self.logger.warning(f"Invalid color hex: {color_hex}")

                # Apply alignment if specified
                if 'alignment' in styles:
                    if styles['alignment'] == 'center':
                        paragraph.alignment = PP_ALIGN.CENTER
                    elif styles['alignment'] == 'right':
                        paragraph.alignment = PP_ALIGN.RIGHT

        except Exception as e:
            self.logger.error(f"Failed to apply text styles to {section_name}: {e}")

    def _replace_placeholder_with_image(self, slide, placeholder, pil_image: Image.Image):
        """Enhanced image replacement with better error handling."""
        try:
            # Get placeholder position and size
            left, top, width, height = placeholder.left, placeholder.top, placeholder.width, placeholder.height

            # Remove the old shape
            try:
                slide.shapes._spTree.remove(placeholder._element)
            except Exception as e:
                self.logger.debug(f"Could not remove placeholder element: {e}")

            # Ensure image is in compatible format
            if pil_image.mode in ('RGBA', 'LA', 'P'):
                pil_image = pil_image.convert('RGB')

            # Save and add image
            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                pil_image.save(tmp, format="PNG")
                path = tmp.name

            pic = slide.shapes.add_picture(path, left, top, width, height)

            # Clean up temporary file
            try:
                Path(path).unlink(missing_ok=True)
            except Exception as e:
                self.logger.debug(f"Could not delete temporary file {path}: {e}")

            self.logger.debug("Successfully replaced placeholder with image")
            return pic

        except Exception as e:
            self.logger.error(f"Failed to replace placeholder with image: {e}")
            return None
