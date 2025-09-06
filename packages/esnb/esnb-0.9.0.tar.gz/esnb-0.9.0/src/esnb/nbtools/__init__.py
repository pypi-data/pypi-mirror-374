import string

import matplotlib.pyplot as plt
import numpy as np

from io import BytesIO
from matplotlib import rcParams
from matplotlib.colors import BoundaryNorm, ListedColormap
from matplotlib.transforms import Bbox

from . import xr_stats

try:
    import momgrid
    from pptx import Presentation
    from pptx.util import Inches
except Exception:
    pass

__all__ = ["xr_stats"]


def gen_levs_and_cmap(start, end, delta, cmap="RdBu_r"):
    """Generates a difference colormap centered on white"""
    boundaries = np.arange(start, end, delta)
    levels = (boundaries[0:-1] + boundaries[1:]) / 2.0
    base_cmap = plt.get_cmap(cmap)
    colors = base_cmap(np.linspace(0, 1, len(levels)))
    colors[[int(len(colors) / 2) - 1]] = [1, 1, 1, 1]
    colors[[int(len(colors) / 2)]] = [1, 1, 1, 1]
    cmap = ListedColormap(colors)
    norm = BoundaryNorm(boundaries, cmap.N, clip=True)
    return (cmap, norm, boundaries)


def set_annotaions(ax, expName="", starttime=None, endtime=None):
    _ = ax.set_xticks([])
    _ = ax.set_yticks([])
    _ = ax.text(
        0.0,
        1.06,
        "SST Bias Relative to NOAA OISSTv2 (1993-2017)",
        weight="bold",
        fontsize=12,
        transform=ax.transAxes,
    )
    _ = ax.text(0.0, 1.02, expName, style="italic", fontsize=10, transform=ax.transAxes)
    _ = ax.text(
        1.0, 1.05, str(starttime.values), ha="right", fontsize=8, transform=ax.transAxes
    )
    _ = ax.text(
        1.0, 1.02, str(endtime.values), ha="right", fontsize=8, transform=ax.transAxes
    )


def add_stats_box(ax, stats_str, x=0.015, y=0.8):
    # Adding the text box annotation
    props = dict(
        boxstyle="round,pad=0.3", edgecolor="black", linewidth=1.5, facecolor="white"
    )
    ax.text(
        x,
        y,
        stats_str,
        transform=ax.transAxes,
        fontsize=8,
        verticalalignment="top",
        bbox=props,
    )


def calculate_stats(model, obs, areacello):
    diff = model - obs
    stats = {}
    stats["min"] = float(diff.min())
    stats["max"] = float(diff.max())
    stats = {**stats, **momgrid.xr_stats.xr_stats_2d(model, obs, areacello, fmt="dict")}
    # Limit to 4 significant digits
    stats = {k: f"{v:.4g}" for k, v in stats.items()}
    # Stats string
    stats_str = str("\n").join([f"{k} = {v}" for k, v in stats.items()])
    return (stats, stats_str)


def add_colorbar(fig, cb, boundaries, limits=[0.16, 0.06, 0.7, 0.03], label=None):
    cbar_ax = fig.add_axes(limits)
    fig.colorbar(
        cb,
        cax=cbar_ax,
        orientation="horizontal",
        extend="both",
        ticks=boundaries[::4],
        label=label,
    )


def setup_plots(dpi=150):
    # Setup figure dpi
    rcParams["figure.dpi"] = dpi

    # Embed font and text in PDF if saved
    rcParams["pdf.fonttype"] = 42.0

    # Modify default font sizes
    plt.rc("font", size=4.5)  # controls default text sizes
    plt.rc("axes", titlesize=7)  # fontsize of the axes title
    plt.rc("axes", labelsize=6)  # fontsize of the x and y labels
    plt.rc("xtick", labelsize=6)  # fontsize of the tick labels
    plt.rc("ytick", labelsize=6)  # fontsize of the tick labels
    plt.rc("legend", fontsize=4.5)  # legend fontsize
    plt.rc("figure", titlesize=10)  # fontsize of the figure title


SINGLE_COLUMN = 3.35
DOUBLE_COLUMN = 5.41
FULL_PAGE = 6.69


def get_figsize_subplots(nexps):
    if nexps == 1:
        figsize = (SINGLE_COLUMN, SINGLE_COLUMN)
        subplot = (1, 1)
    elif nexps == 2:
        figsize = (DOUBLE_COLUMN, SINGLE_COLUMN)
        subplot = (1, 2)
    elif nexps == 3:
        figsize = (FULL_PAGE, SINGLE_COLUMN)
        subplot = (1, 3)
    elif nexps == 4:
        figsize = (DOUBLE_COLUMN, DOUBLE_COLUMN)
        subplot = (2, 2)
    else:
        nrows = int(np.ceil(nexps / 3.0))
        figsize = (FULL_PAGE, nrows * SINGLE_COLUMN)
        subplot = (nrows, 3)

    return (figsize, subplot)


def bottom_colorbar(fig, cb, **kwargs):
    axs = [
        ax
        for ax in fig.get_axes()
        if ax.get_visible() and not ax.get_label().startswith("colorbar")
    ]

    fig.canvas.draw()
    bbox = axs[0].get_tightbbox(fig.canvas.get_renderer())

    # Compute union of all visible axes' tight bounding boxes
    bbox = axs[0].get_tightbbox(fig.canvas.get_renderer())
    for ax in axs[1:]:
        bbox = Bbox.union([bbox, ax.get_tightbbox(fig.canvas.get_renderer())])

    bbox_fig = bbox.transformed(fig.transFigure.inverted())

    cbar_height = 0.02
    cbar_pad = 0.05
    cbar_bottom = bbox_fig.y0 - cbar_pad
    cbar_left = bbox_fig.x0
    cbar_width = bbox_fig.width

    cbar_ax = fig.add_axes([cbar_left, cbar_bottom, cbar_width, cbar_height])
    cbar = fig.colorbar(cb, cax=cbar_ax, **kwargs)
    return cbar


def panel_letters(axes, x=-0.1, y=1.12, **kwargs):
    labels = [f"{x}." for x in list(string.ascii_lowercase)]
    if kwargs == {}:
        kwargs = {"weight": "bold", "ha": "left", "fontsize": 7}
    for n, ax in enumerate(axes):
        ax.text(x, y, labels[n], transform=ax.transAxes, **kwargs)


def save_pptx(figs, filename, dpi=200):
    figs = [figs] if not isinstance(figs, list) else figs

    prs = Presentation()
    # 16x9 layout below
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    # Loop through figures
    for fig in figs:
        image_stream = BytesIO()
        fig.savefig(image_stream, format="png", bbox_inches="tight", dpi=dpi)
        image_stream.seek(0)

        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(
            image_stream, Inches(0.5), Inches(0.5), width=Inches(12.33)
        )  # 0.5-inch margin

    # Save the presentation
    prs.save(filename)
